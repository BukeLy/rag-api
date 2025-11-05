"""
Document Complexity Analyzer

基于 v2.0 智能 Parser 选择方案的复杂度评分算法

功能特性：
- 多维度复杂度评分：图片、表格、布局、字体、文本密度
- 关键决策因子：表格行数、中文密度
- 支持 PDF、Office 文档、图片
- 采样分析（大文件优化）
"""

import os
import re
from typing import Dict, Any, Optional
from pathlib import Path
from dataclasses import dataclass

# PDF 分析
import fitz  # PyMuPDF

# Office 文档分析
try:
    import pdfplumber
    from docx import Document as DocxDocument
    from pptx import Presentation
    import openpyxl
except ImportError:
    pdfplumber = None
    DocxDocument = None
    Presentation = None
    openpyxl = None

from src.logger import logger


@dataclass
class DocumentFeatures:
    """文档特征"""
    # 基础信息
    page_count: int
    file_size_kb: float

    # 图片特征
    avg_image_count_per_page: float
    total_image_count: int

    # 表格特征
    avg_table_count_per_page: float
    total_table_count: int
    avg_table_row_count: float  # 新增：表格平均行数（关键决策因子）

    # 布局特征
    has_complex_layout: bool
    font_variety: int

    # 文本特征
    text_density: float  # 文本密度（0-1）
    low_text_density: bool  # 是否低文本密度

    # 中文特征（新增）
    chinese_char_count: int
    total_char_count: int
    chinese_char_ratio: float  # 中文字符占比（关键决策因子）


class DocumentComplexityAnalyzer:
    """文档复杂度分析器（基于 v2.0 实测调优）"""

    # 环境变量配置
    SIMPLE_THRESHOLD = int(os.getenv("COMPLEXITY_SIMPLE_THRESHOLD", "20"))
    MEDIUM_TABLE_THRESHOLD = int(os.getenv("COMPLEXITY_MEDIUM_TABLE_THRESHOLD", "40"))
    COMPLEX_SINGLE_PAGE_THRESHOLD = int(os.getenv("COMPLEXITY_COMPLEX_SINGLE_PAGE_THRESHOLD", "60"))

    # 采样配置（大文件优化）
    MAX_SAMPLE_PAGES = 20
    SAMPLE_RATIO = 0.1  # 采样 10% 页面

    def __init__(self):
        """初始化分析器"""
        logger.info("DocumentComplexityAnalyzer initialized")

    def analyze_complexity(self, file_path: str) -> int:
        """
        计算文档复杂度评分

        Args:
            file_path: 文件路径

        Returns:
            复杂度评分（0-100+）

        公式（基于 v2.0）：
        score = (
            avg_image_count_per_page * 10 +
            avg_table_count_per_page * 15 +
            has_complex_layout * 20 +
            font_variety * 3 +
            low_text_density * 10 +
            avg_table_row_count_per_page * 1 +  # 新增
            chinese_char_ratio * 10                # 新增
        )
        """
        features = self.get_document_features(file_path)

        score = (
            features.avg_image_count_per_page * 10 +
            features.avg_table_count_per_page * 15 +
            (20 if features.has_complex_layout else 0) +
            features.font_variety * 3 +
            (10 if features.low_text_density else 0) +
            features.avg_table_row_count * 1 +
            features.chinese_char_ratio * 10
        )

        logger.info(
            f"Complexity score: {score:.1f} "
            f"(images={features.avg_image_count_per_page:.1f}, "
            f"tables={features.avg_table_count_per_page:.1f}, "
            f"table_rows={features.avg_table_row_count:.1f}, "
            f"chinese={features.chinese_char_ratio:.2%})"
        )

        return int(score)

    def get_document_features(self, file_path: str) -> DocumentFeatures:
        """
        提取文档特征

        Args:
            file_path: 文件路径

        Returns:
            DocumentFeatures 对象
        """
        file_path = Path(file_path)
        ext = file_path.suffix.lower()

        if ext == ".pdf":
            return self._analyze_pdf(file_path)
        elif ext in [".doc", ".docx"]:
            return self._analyze_docx(file_path)
        elif ext in [".ppt", ".pptx"]:
            return self._analyze_pptx(file_path)
        elif ext in [".xls", ".xlsx"]:
            return self._analyze_xlsx(file_path)
        elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".gif"]:
            return self._analyze_image(file_path)
        else:
            # 未知格式，返回默认特征
            logger.warning(f"Unsupported file format: {ext}, using default features")
            return self._get_default_features(file_path)

    def _analyze_pdf(self, file_path: Path) -> DocumentFeatures:
        """
        分析 PDF 文档

        使用 PyMuPDF + pdfplumber 双引擎：
        - PyMuPDF：快速提取图片、字体、文本
        - pdfplumber：准确提取表格
        """
        doc = fitz.open(file_path)
        page_count = len(doc)

        # 采样策略（大文件优化）
        sample_pages = self._get_sample_pages(page_count)
        sample_count = len(sample_pages)

        total_images = 0
        total_tables = 0
        total_table_rows = 0
        fonts_set = set()
        total_text_len = 0
        total_page_area = 0
        chinese_chars = 0
        total_chars = 0

        # 使用 PyMuPDF 快速提取基础信息
        for page_num in sample_pages:
            page = doc[page_num]

            # 图片数量
            images = page.get_images()
            total_images += len(images)

            # 字体种类
            fonts = page.get_fonts()
            for font in fonts:
                fonts_set.add(font[3])  # font name

            # 文本长度
            text = page.get_text()
            total_text_len += len(text)

            # 页面面积（用于计算文本密度）
            rect = page.rect
            total_page_area += rect.width * rect.height

            # 中文字符统计
            chinese_chars += len(re.findall(r'[\u4e00-\u9fff]', text))
            total_chars += len(text.strip())

        doc.close()

        # 使用 pdfplumber 提取表格（更准确）
        if pdfplumber:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page_num in sample_pages:
                        if page_num < len(pdf.pages):
                            page = pdf.pages[page_num]
                            tables = page.extract_tables()
                            if tables:
                                total_tables += len(tables)
                                # 统计表格行数
                                for table in tables:
                                    if table:
                                        total_table_rows += len(table)
            except Exception as e:
                logger.warning(f"pdfplumber failed, skipping table extraction: {e}")

        # 计算平均值
        avg_images = total_images / sample_count if sample_count > 0 else 0
        avg_tables = total_tables / sample_count if sample_count > 0 else 0
        avg_table_rows = total_table_rows / total_tables if total_tables > 0 else 0

        # 文本密度（文本长度 / 页面面积）
        text_density = total_text_len / total_page_area if total_page_area > 0 else 0
        low_text_density = text_density < 0.001  # 经验阈值

        # 复杂布局判断（多列布局、复杂排版）
        # 简化判断：字体种类 >5 认为复杂
        has_complex_layout = len(fonts_set) > 5

        # 中文字符占比
        chinese_ratio = chinese_chars / total_chars if total_chars > 0 else 0

        return DocumentFeatures(
            page_count=page_count,
            file_size_kb=file_path.stat().st_size / 1024,
            avg_image_count_per_page=avg_images,
            total_image_count=total_images,
            avg_table_count_per_page=avg_tables,
            total_table_count=total_tables,
            avg_table_row_count=avg_table_rows,
            has_complex_layout=has_complex_layout,
            font_variety=len(fonts_set),
            text_density=text_density,
            low_text_density=low_text_density,
            chinese_char_count=chinese_chars,
            total_char_count=total_chars,
            chinese_char_ratio=chinese_ratio,
        )

    def _analyze_docx(self, file_path: Path) -> DocumentFeatures:
        """分析 Word 文档"""
        if not DocxDocument:
            logger.warning("python-docx not installed, returning default features")
            return self._get_default_features(file_path)

        doc = DocxDocument(file_path)

        # 统计图片（通过 relationships）
        total_images = len([r for r in doc.part.rels.values() if "image" in r.target_ref])

        # 统计表格
        total_tables = len(doc.tables)
        total_table_rows = sum(len(table.rows) for table in doc.tables)
        avg_table_rows = total_table_rows / total_tables if total_tables > 0 else 0

        # 统计段落和文本
        page_count = max(1, len(doc.sections))  # 近似页数（section 数量）
        text = "\n".join([p.text for p in doc.paragraphs])
        chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', text))
        total_chars = len(text.strip())
        chinese_ratio = chinese_chars / total_chars if total_chars > 0 else 0

        return DocumentFeatures(
            page_count=page_count,
            file_size_kb=file_path.stat().st_size / 1024,
            avg_image_count_per_page=total_images / page_count,
            total_image_count=total_images,
            avg_table_count_per_page=total_tables / page_count,
            total_table_count=total_tables,
            avg_table_row_count=avg_table_rows,
            has_complex_layout=False,  # Word 文档通常布局简单
            font_variety=3,  # 假设 3 种字体
            text_density=0.5,  # 假设中等密度
            low_text_density=False,
            chinese_char_count=chinese_chars,
            total_char_count=total_chars,
            chinese_char_ratio=chinese_ratio,
        )

    def _analyze_pptx(self, file_path: Path) -> DocumentFeatures:
        """分析 PowerPoint 文档"""
        if not Presentation:
            logger.warning("python-pptx not installed, returning default features")
            return self._get_default_features(file_path)

        prs = Presentation(file_path)
        page_count = len(prs.slides)

        total_images = 0
        total_tables = 0
        total_text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                # 图片
                if shape.shape_type == 13:  # PICTURE
                    total_images += 1
                # 表格
                elif shape.shape_type == 19:  # TABLE
                    total_tables += 1
                # 文本
                if hasattr(shape, "text"):
                    total_text += shape.text

        chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', total_text))
        total_chars = len(total_text.strip())
        chinese_ratio = chinese_chars / total_chars if total_chars > 0 else 0

        return DocumentFeatures(
            page_count=page_count,
            file_size_kb=file_path.stat().st_size / 1024,
            avg_image_count_per_page=total_images / page_count if page_count > 0 else 0,
            total_image_count=total_images,
            avg_table_count_per_page=total_tables / page_count if page_count > 0 else 0,
            total_table_count=total_tables,
            avg_table_row_count=5,  # PPT 表格通常较小
            has_complex_layout=True,  # PPT 布局通常复杂
            font_variety=5,  # 假设 5 种字体
            text_density=0.3,  # PPT 文本密度通常较低
            low_text_density=True,
            chinese_char_count=chinese_chars,
            total_char_count=total_chars,
            chinese_char_ratio=chinese_ratio,
        )

    def _analyze_xlsx(self, file_path: Path) -> DocumentFeatures:
        """分析 Excel 文档"""
        if not openpyxl:
            logger.warning("openpyxl not installed, returning default features")
            return self._get_default_features(file_path)

        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheet_count = len(wb.sheetnames)

        total_tables = sheet_count  # 每个 sheet 算一个表格
        total_rows = 0
        total_text = ""

        for sheet in wb.worksheets:
            total_rows += sheet.max_row
            # 提取文本（仅采样前 100 行）
            for row in sheet.iter_rows(max_row=100, values_only=True):
                for cell in row:
                    if cell:
                        total_text += str(cell)

        avg_table_rows = total_rows / total_tables if total_tables > 0 else 0

        chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', total_text))
        total_chars = len(total_text.strip())
        chinese_ratio = chinese_chars / total_chars if total_chars > 0 else 0

        return DocumentFeatures(
            page_count=sheet_count,
            file_size_kb=file_path.stat().st_size / 1024,
            avg_image_count_per_page=0,  # Excel 图片暂不统计
            total_image_count=0,
            avg_table_count_per_page=1,  # 每页一个表格
            total_table_count=total_tables,
            avg_table_row_count=avg_table_rows,
            has_complex_layout=False,
            font_variety=2,
            text_density=0.6,
            low_text_density=False,
            chinese_char_count=chinese_chars,
            total_char_count=total_chars,
            chinese_char_ratio=chinese_ratio,
        )

    def _analyze_image(self, file_path: Path) -> DocumentFeatures:
        """分析图片文件"""
        # 图片文件：1 页、1 个图片、0 个表格
        return DocumentFeatures(
            page_count=1,
            file_size_kb=file_path.stat().st_size / 1024,
            avg_image_count_per_page=1,
            total_image_count=1,
            avg_table_count_per_page=0,
            total_table_count=0,
            avg_table_row_count=0,
            has_complex_layout=False,
            font_variety=0,
            text_density=0,
            low_text_density=True,
            chinese_char_count=0,
            total_char_count=0,
            chinese_char_ratio=0,
        )

    def _get_default_features(self, file_path: Path) -> DocumentFeatures:
        """返回默认特征（用于未知格式）"""
        return DocumentFeatures(
            page_count=1,
            file_size_kb=file_path.stat().st_size / 1024,
            avg_image_count_per_page=0,
            total_image_count=0,
            avg_table_count_per_page=0,
            total_table_count=0,
            avg_table_row_count=0,
            has_complex_layout=False,
            font_variety=0,
            text_density=0.5,
            low_text_density=False,
            chinese_char_count=0,
            total_char_count=0,
            chinese_char_ratio=0,
        )

    def _get_sample_pages(self, page_count: int) -> list[int]:
        """
        获取采样页码列表

        大文件优化：
        - 页数 ≤ 20：全部分析
        - 页数 > 20：采样 10%，最多 20 页
        """
        if page_count <= self.MAX_SAMPLE_PAGES:
            return list(range(page_count))

        # 采样策略：均匀采样
        sample_count = min(
            self.MAX_SAMPLE_PAGES,
            max(1, int(page_count * self.SAMPLE_RATIO))
        )

        step = page_count // sample_count
        return [i * step for i in range(sample_count)]
